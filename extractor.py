#!/usr/bin/env python3
"""
PPTX/PDF 콘텐츠 추출 모듈
inputdata/의 파일들에서 텍스트와 구조를 추출합니다.

XML 기반 추출로 모든 텍스트를 완전히 추출합니다.
"""

import re
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
import pdfplumber


# XML 네임스페이스 정의
NAMESPACES = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
}


def extract_text_from_xml(xml_content: bytes) -> list:
    """
    XML 콘텐츠에서 모든 텍스트를 추출합니다.
    <a:t> 태그의 모든 텍스트를 추출합니다.
    """
    texts = []
    try:
        root = ET.fromstring(xml_content)

        # 모든 <a:t> 태그 찾기 (텍스트 요소)
        for elem in root.iter():
            # 네임스페이스 포함된 태그명 확인
            if elem.tag.endswith('}t') and 'drawingml' in elem.tag:
                if elem.text:
                    texts.append(elem.text)
    except ET.ParseError:
        pass

    return texts


def extract_paragraphs_from_xml(xml_content: bytes) -> list:
    """
    XML에서 문단 단위로 텍스트를 추출합니다.
    <a:p> (paragraph) 단위로 텍스트를 그룹화합니다.
    """
    paragraphs = []
    try:
        root = ET.fromstring(xml_content)

        # drawingml 네임스페이스의 paragraph 태그 찾기
        for elem in root.iter():
            if elem.tag.endswith('}p') and 'drawingml' in elem.tag:
                # 이 문단 내의 모든 텍스트 수집
                para_texts = []
                for child in elem.iter():
                    if child.tag.endswith('}t') and 'drawingml' in child.tag:
                        if child.text:
                            para_texts.append(child.text)

                if para_texts:
                    paragraph = ''.join(para_texts).strip()
                    if paragraph:
                        paragraphs.append(paragraph)
    except ET.ParseError:
        pass

    return paragraphs


def extract_shapes_from_xml(xml_content: bytes) -> list:
    """
    XML에서 shape(도형) 단위로 텍스트를 추출합니다.
    각 shape 내의 텍스트를 하나의 블록으로 그룹화합니다.
    """
    shapes = []
    try:
        root = ET.fromstring(xml_content)

        # sp (shape), pic (picture), graphicFrame 등에서 텍스트 추출
        shape_tags = ['}sp', '}graphicFrame', '}cxnSp']

        for elem in root.iter():
            is_shape = any(elem.tag.endswith(tag) for tag in shape_tags)
            if is_shape:
                # shape 내의 모든 문단 텍스트 수집
                shape_paragraphs = []
                for child in elem.iter():
                    if child.tag.endswith('}p') and 'drawingml' in child.tag:
                        para_texts = []
                        for t_elem in child.iter():
                            if t_elem.tag.endswith('}t') and 'drawingml' in t_elem.tag:
                                if t_elem.text:
                                    para_texts.append(t_elem.text)
                        if para_texts:
                            shape_paragraphs.append(''.join(para_texts))

                if shape_paragraphs:
                    shape_text = '\n'.join(shape_paragraphs).strip()
                    if shape_text:
                        shapes.append(shape_text)
    except ET.ParseError:
        pass

    return shapes


def extract_pptx_xml(file_path: str) -> dict:
    """
    PPTX 파일을 ZIP으로 열어 XML에서 직접 모든 텍스트를 추출합니다.

    Returns:
        {
            "source": "파일명",
            "slide_count": 슬라이드 수,
            "slides": [
                {"page": 1, "texts": ["텍스트1", "텍스트2", ...]}
            ]
        }
    """
    slides_data = []

    with zipfile.ZipFile(file_path, 'r') as zf:
        # 슬라이드 파일 목록 가져오기
        slide_files = sorted([
            f for f in zf.namelist()
            if f.startswith('ppt/slides/slide') and f.endswith('.xml')
        ], key=lambda x: int(re.search(r'slide(\d+)', x).group(1)))

        for idx, slide_file in enumerate(slide_files):
            xml_content = zf.read(slide_file)

            # shape 단위로 텍스트 추출 (가장 구조화된 방식)
            texts = extract_shapes_from_xml(xml_content)

            # shape에서 추출 안 된 경우 문단 단위로 시도
            if not texts:
                texts = extract_paragraphs_from_xml(xml_content)

            slides_data.append({
                "page": idx + 1,
                "texts": texts
            })

    return {
        "source": Path(file_path).name,
        "slide_count": len(slides_data),
        "slides": slides_data
    }


def extract_pptx(file_path: str) -> dict:
    """
    PPTX 파일에서 슬라이드별 텍스트를 추출합니다.
    XML 기반 추출을 우선 시도하고, 실패 시 python-pptx 사용.

    Returns:
        {
            "source": "파일명",
            "slide_count": 슬라이드 수,
            "slides": [
                {"page": 1, "texts": ["텍스트1", "텍스트2", ...]}
            ]
        }
    """
    # XML 기반 추출 시도
    try:
        result = extract_pptx_xml(file_path)
        # 텍스트가 추출되었는지 확인
        has_text = any(slide.get("texts") for slide in result.get("slides", []))
        if has_text:
            return result
    except Exception as e:
        print(f"  XML 추출 실패, python-pptx로 시도: {e}")

    # Fallback: python-pptx 사용
    prs = Presentation(file_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)

        slides_data.append({
            "page": idx + 1,
            "texts": texts
        })

    return {
        "source": Path(file_path).name,
        "slide_count": len(slides_data),
        "slides": slides_data
    }


def extract_pdf(file_path: str) -> dict:
    """
    PDF 파일에서 페이지별 텍스트를 추출합니다.

    Returns:
        {
            "source": "파일명",
            "page_count": 페이지 수,
            "pages": [
                {"page": 1, "text": "페이지 텍스트"}
            ]
        }
    """
    pages_data = []

    with pdfplumber.open(file_path) as pdf:
        for idx, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            pages_data.append({
                "page": idx + 1,
                "text": text.strip()
            })

    return {
        "source": Path(file_path).name,
        "page_count": len(pages_data),
        "pages": pages_data
    }


def detect_document_type(filename: str) -> str:
    """파일명에서 문서 유형을 추출합니다."""
    filename_lower = filename.lower()

    if "제안서" in filename:
        return "제안서"
    elif "결과보고서" in filename:
        return "결과보고서"
    elif "기획안" in filename:
        return "기획안"
    elif "백서" in filename:
        return "백서"
    elif "컨설팅" in filename or "보고서" in filename:
        return "컨설팅보고서"
    elif "중장기" in filename or "계획" in filename:
        return "중장기계획"
    else:
        return "기타"


def extract_chapter_pattern(texts: list) -> list:
    """
    텍스트에서 챕터 패턴을 추출합니다.
    예: "01. 사업개요", "1. 소개", "I. 개요" 등
    """
    chapter_patterns = [
        r"^(\d{2})\.\s*(.+)$",      # 01. 사업개요
        r"^(\d+)\.\s*(.+)$",        # 1. 사업개요
        r"^([IVX]+)\.\s*(.+)$",     # I. 사업개요
        r"^(Chapter\s*\d+)[\.:]\s*(.+)$",  # Chapter 1: 소개
    ]

    chapters = []
    for text in texts:
        for line in text.split('\n'):
            line = line.strip()
            for pattern in chapter_patterns:
                match = re.match(pattern, line)
                if match:
                    chapters.append(line)
                    break

    return chapters


def is_noise_text(text: str) -> bool:
    """노이즈 텍스트인지 판별합니다."""
    if not text or len(text.strip()) < 3:
        return True

    noise_patterns = [
        r"^[‹›\d#]+$",                    # 페이지 번호 (‹#›, 1, 2, 등)
        r"^ⓒ\s*\d{4}",                    # 저작권 문구
        r"^©\s*\d{4}",                    # 저작권 문구
        r"All rights reserved",            # 저작권 문구
        r"^\d+$",                          # 숫자만
        r"^[IVXivx]+$",                   # 로마 숫자만
        r"^[\s\n]+$",                      # 공백만
        r"^(Click to|클릭하여)",           # 플레이스홀더
        r"^(Title|Subtitle|제목)",         # 플레이스홀더
    ]

    text_stripped = text.strip()
    for pattern in noise_patterns:
        if re.search(pattern, text_stripped, re.IGNORECASE):
            return True

    return False


def filter_meaningful_texts(texts: list) -> list:
    """의미있는 텍스트만 필터링합니다."""
    return [t for t in texts if not is_noise_text(t)]


def analyze_slide_structure(slide_data: dict) -> dict:
    """
    슬라이드 텍스트에서 chapter/title/lead 구조를 추정합니다.
    노이즈 텍스트를 제외하고 의미있는 콘텐츠만 추출합니다.
    """
    raw_texts = slide_data.get("texts", [])
    texts = filter_meaningful_texts(raw_texts)

    result = {
        "page": slide_data.get("page"),
        "chapter": None,
        "title": None,
        "lead": None,
        "raw_texts": raw_texts
    }

    if not texts:
        return result

    # 텍스트 길이순으로 정렬 (짧은 것이 title, 긴 것이 lead일 가능성)
    # 하지만 순서도 중요하므로 처음 몇 개만 분석

    # 첫 번째 텍스트 분석
    first_text = texts[0]

    # 챕터 패턴 확인 (01. 섹션명, Ⅰ. 섹션명 등)
    chapter_patterns = [
        r"^(\d{1,2})\.\s*(.+)$",           # 1. 또는 01. 사업개요
        r"^([IVXⅠⅡⅢⅣⅤⅥⅦⅧⅨⅩivx]+)[\.\s]+(.+)$",  # Ⅰ 성과요약, I. 소개
    ]

    chapter_found = False
    for pattern in chapter_patterns:
        match = re.match(pattern, first_text)
        if match:
            result["chapter"] = first_text
            chapter_found = True
            break

    if chapter_found:
        # 챕터 다음 텍스트들에서 title, lead 추출
        remaining = texts[1:]
        if remaining:
            # 가장 짧고 의미있는 것을 title로
            result["title"] = remaining[0]
        if len(remaining) > 1:
            # 그 다음 긴 텍스트를 lead로
            result["lead"] = remaining[1]
    else:
        # 챕터가 없는 경우: 짧은 것 title, 긴 것 lead
        if len(texts) >= 1:
            # 길이가 적당한 것을 title로 (5~50자)
            for t in texts:
                if 5 <= len(t) <= 50:
                    result["title"] = t
                    break
            if not result["title"]:
                result["title"] = texts[0]

        if len(texts) >= 2:
            # title 외의 가장 긴 텍스트를 lead로
            for t in texts:
                if t != result["title"] and len(t) > 10:
                    result["lead"] = t
                    break

    return result


def extract_all_from_directory(input_dir: str) -> list:
    """
    디렉토리의 모든 PPTX/PDF 파일을 추출합니다.

    Returns:
        [
            {
                "source": "파일명",
                "type": "문서유형",
                "format": "pptx" | "pdf",
                "content": 추출된 데이터
            }
        ]
    """
    input_path = Path(input_dir)
    results = []

    # PPTX 파일 처리
    for pptx_file in input_path.glob("*.pptx"):
        if pptx_file.name.startswith(("~$", "tem", "output")):
            continue

        try:
            print(f"추출 중: {pptx_file.name}")
            content = extract_pptx(str(pptx_file))
            results.append({
                "source": pptx_file.name,
                "type": detect_document_type(pptx_file.name),
                "format": "pptx",
                "content": content
            })
        except Exception as e:
            print(f"  오류: {e}")

    # PDF 파일 처리
    for pdf_file in input_path.glob("*.pdf"):
        try:
            print(f"추출 중: {pdf_file.name}")
            content = extract_pdf(str(pdf_file))
            results.append({
                "source": pdf_file.name,
                "type": detect_document_type(pdf_file.name),
                "format": "pdf",
                "content": content
            })
        except Exception as e:
            print(f"  오류: {e}")

    return results


if __name__ == "__main__":
    # 테스트 실행
    import json

    input_dir = "inputdata"
    results = extract_all_from_directory(input_dir)

    print(f"\n총 {len(results)}개 파일 추출 완료")
    for r in results:
        if r["format"] == "pptx":
            print(f"  - {r['source']}: {r['content']['slide_count']}슬라이드 ({r['type']})")
        else:
            print(f"  - {r['source']}: {r['content']['page_count']}페이지 ({r['type']})")
