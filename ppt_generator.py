#!/usr/bin/env python3
"""
PPTX 자동 생성기
템플릿의 placeholder(Chapter, Title, Lead)를 JSON 데이터로 채워서 슬라이드를 생성합니다.

텍스트 매칭 방식: 템플릿에서 "Chapter", "Title", "Lead" 텍스트를 찾아서 대체합니다.
"""

import json
import copy
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


def load_json_data(json_path: str) -> dict:
    """JSON 파일에서 슬라이드 데이터를 로드합니다."""
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def duplicate_slide(prs: Presentation, index: int):
    """슬라이드를 복제합니다."""
    template = prs.slides[index]
    slide_layout = template.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # 기존 shapes 제거
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # 템플릿 shapes 복사
    for shape in template.shapes:
        new_el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide


def apply_font_style(run, font_config: dict):
    """run에 폰트 스타일을 적용합니다."""
    if not font_config:
        return

    if 'name' in font_config:
        font_name = font_config['name']
        run.font.name = font_name

        # 한글(동아시아) 폰트도 함께 설정
        from pptx.oxml.ns import qn
        rPr = run._r.get_or_add_rPr()

        # a:ea (East Asian) 폰트 설정
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            from lxml import etree
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font_name)

    if 'size' in font_config:
        run.font.size = Pt(font_config['size'])
    if 'bold' in font_config:
        run.font.bold = font_config['bold']
    if 'italic' in font_config:
        run.font.italic = font_config['italic']
    if 'color' in font_config:
        # HEX 색상 지원 (예: "#FF0000")
        color = font_config['color'].lstrip('#')
        run.font.color.rgb = RGBColor(
            int(color[0:2], 16),
            int(color[2:4], 16),
            int(color[4:6], 16)
        )


def find_and_replace_text(slide, placeholder: str, new_text: str, font_config: dict = None):
    """
    슬라이드에서 placeholder 텍스트를 찾아 대체합니다.
    템플릿에 상관없이 텍스트 내용으로 매칭합니다.
    """
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue

        # 텍스트 내용으로 매칭 (대소문자 무시)
        if placeholder.lower() in shape.text_frame.text.lower():
            # 첫 번째 paragraph에서 처리
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]

                # run이 있으면 첫 번째 run에 텍스트 설정, 나머지 run 삭제
                if para.runs:
                    # 첫 번째 run에 텍스트와 스타일 적용
                    para.runs[0].text = new_text
                    apply_font_style(para.runs[0], font_config)

                    # 나머지 run들의 텍스트 비우기
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    # run이 없으면 paragraph 텍스트 직접 설정
                    para.text = new_text

                return True
    return False


def update_slide(slide, slide_data: dict, font_settings: dict = None):
    """슬라이드의 모든 placeholder를 업데이트합니다."""
    font_settings = font_settings or {}

    # Subtitle을 Title보다 먼저 처리 (Title이 Subtitle에 포함되어 있으므로)
    find_and_replace_text(slide, "Subtitle", slide_data.get('subtitle', ''), font_settings.get('subtitle'))
    find_and_replace_text(slide, "Title", slide_data.get('title', ''), font_settings.get('title'))
    find_and_replace_text(slide, "Chapter", slide_data.get('chapter', ''), font_settings.get('chapter'))
    find_and_replace_text(slide, "Lead", slide_data.get('lead', ''), font_settings.get('lead'))
    find_and_replace_text(slide, "Content", slide_data.get('content', ''), font_settings.get('content'))
    find_and_replace_text(slide, "Manager", slide_data.get('manager', ''), font_settings.get('manager'))


def generate_pptx(template_path: str, json_path: str, output_path: str):
    """JSON 데이터를 기반으로 PPTX를 생성합니다."""
    data = load_json_data(json_path)
    slides_data = data.get('slides', [])
    font_settings = data.get('font_settings', {})  # 선택적 폰트 설정

    if not slides_data:
        print("JSON에 슬라이드 데이터가 없습니다.")
        return

    prs = Presentation(template_path)

    # 필요한 만큼 슬라이드 복제
    for _ in range(len(slides_data) - 1):
        duplicate_slide(prs, 0)

    # 모든 슬라이드에 데이터 적용
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides[i]
        update_slide(slide, slide_data, font_settings)
        print(f"Page {slide_data['page']}: {slide_data['title']}")

    prs.save(output_path)
    print(f"\n생성 완료: {output_path}")
    print(f"총 {len(slides_data)}개의 슬라이드가 생성되었습니다.")


def analyze_template(template_path: str):
    """템플릿의 구조를 분석하여 출력합니다."""
    prs = Presentation(template_path)
    print(f"템플릿 분석: {template_path}")
    print(f"총 슬라이드: {len(prs.slides)}개\n")

    for i, slide in enumerate(prs.slides):
        print(f"=== Slide {i + 1} ===")
        for shape in slide.shapes:
            print(f"  Shape: {shape.name}")
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    print(f"    Text: '{text}'")
        print()


if __name__ == "__main__":
    import sys

    if len(sys.argv) >= 2 and sys.argv[1] == "--analyze":
        # 템플릿 분석 모드
        template = sys.argv[2] if len(sys.argv) >= 3 else "tem.pptx"
        analyze_template(template)
    else:
        # 생성 모드
        template = "tem.pptx"
        json_file = "slides_data.json"
        output = "output.pptx"

        if len(sys.argv) >= 2:
            json_file = sys.argv[1]
        if len(sys.argv) >= 3:
            output = sys.argv[2]
        if len(sys.argv) >= 4:
            template = sys.argv[3]

        generate_pptx(template, json_file, output)
